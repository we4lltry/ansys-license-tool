"""
create_template.py
------------------
'라이선스_확인서_템플릿.pptx' 를 코드로 자동 생성합니다.
app.py 가 기대하는 placeholder idx 와 테이블 구조를 포함합니다.

실행:
    python create_template.py

결과:
    라이선스_확인서_템플릿.pptx  (이 스크립트와 같은 폴더에 생성)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm
import copy, os

# ── 색상 ────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x23, 0x57)
GOLD   = RGBColor(0xB8, 0x92, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0xF0, 0xF0, 0xF0)

# ── 슬라이드 크기 (16:9) ────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]   # Blank
slide = prs.slides.add_slide(slide_layout)

W = prs.slide_width
H = prs.slide_height


# ═══════════════════════════════════════════════════
#  헬퍼 함수
# ═══════════════════════════════════════════════════
def add_textbox(slide, left, top, width, height,
                text="", fontsize=14, bold=False,
                color=RGBColor(0,0,0), align=PP_ALIGN.LEFT,
                ph_idx=None):
    """일반 텍스트박스 or placeholder 역할 텍스트박스 추가."""
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.color.rgb = color

    # placeholder idx 를 태그에 저장 — app.py 는 shape.placeholder_format.idx 를
    # 사용하므로, 실제 레이아웃 placeholder 를 만들어야 합니다.
    # (텍스트박스로는 ph_idx 가 동작하지 않으므로, 아래 별도 함수 사용)
    return txBox


def add_placeholder_shape(slide, ph_idx, left, top, width, height,
                           default_text="", fontsize=12, bold=False,
                           color=RGBColor(0x33, 0x33, 0x33),
                           align=PP_ALIGN.LEFT):
    """
    python-pptx 에서 직접 placeholder 를 추가하는 공식 API 가 없으므로
    'body' placeholder 를 XML 레벨에서 주입합니다.
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    sp_tree = slide.shapes._spTree

    # ── sp XML 뼈대 ──────────────────────────────────
    sp_xml = f"""
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
      <p:nvSpPr>
        <p:cNvPr id="{100 + ph_idx}" name="ph_{ph_idx}"/>
        <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
        <p:nvPr>
          <p:ph type="body" idx="{ph_idx}"/>
        </p:nvPr>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="{left}" y="{top}"/>
          <a:ext cx="{width}" cy="{height}"/>
        </a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:noFill/>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" lIns="45720" rIns="45720"
                  tIns="36000" bIns="36000" anchor="ctr"/>
        <a:lstStyle/>
        <a:p>
          <a:pPr algn="{'ctr' if align==PP_ALIGN.CENTER else 'l'}"/>
          <a:r>
            <a:rPr lang="ko-KR" sz="{int(fontsize*100)}"
                   b="{'1' if bold else '0'}" dirty="0">
              <a:solidFill>
                <a:srgbClr val="{color[0]:02X}{color[1]:02X}{color[2]:02X}"/>
              </a:solidFill>
            </a:rPr>
            <a:t>{default_text}</a:t>
          </a:r>
        </a:p>
      </p:txBody>
    </p:sp>
    """
    sp_elem = etree.fromstring(sp_xml)
    sp_tree.append(sp_elem)


def set_cell(cell, text, fontsize=11, bold=False,
             fg=RGBColor(0,0,0), bg=None, align=PP_ALIGN.LEFT):
    """테이블 셀 텍스트 + 스타일."""
    from pptx.oxml.ns import qn
    from lxml import etree

    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.color.rgb = fg

    if bg:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", f"{bg[0]:02X}{bg[1]:02X}{bg[2]:02X}")


# ═══════════════════════════════════════════════════
#  배경 (네이비 상단 바)
# ═══════════════════════════════════════════════════
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree

def add_rect(slide, left, top, width, height, fill_color):
    sp = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    sp.line.fill.background()       # 테두리 없음
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    return sp

# 상단 네이비 배경 바
add_rect(slide, 0, 0, W, Cm(2.8), NAVY)

# 하단 골드 라인
add_rect(slide, 0, Cm(2.5), W, Cm(0.18), GOLD)


# ═══════════════════════════════════════════════════
#  제목 (상단 바 안에)
# ═══════════════════════════════════════════════════
title_box = add_textbox(
    slide, Cm(1.5), Cm(0.5), Cm(20), Cm(1.8),
    text="라 이 선 스  확 인 서",
    fontsize=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER
)

# 회사명 (우측 상단)
add_textbox(
    slide, W - Cm(7), Cm(0.7), Cm(6.5), Cm(1.2),
    text="(주) 태성에스엔이",
    fontsize=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT
)


# ═══════════════════════════════════════════════════
#  고객 정보 레이블 (고정 텍스트)
# ═══════════════════════════════════════════════════
LABEL_COLOR = RGBColor(0x1a, 0x23, 0x57)
ROW_H = Cm(0.85)
COL1_L = Cm(1.5)
COL1_W = Cm(3.5)
COL2_L = Cm(5.0)
COL2_W = Cm(7.5)
COL3_L = Cm(12.8)
COL3_W = Cm(2.5)
COL4_L = Cm(15.5)
COL4_W = Cm(5.0)

info_top = Cm(3.2)

labels_left  = ["고   객   명",  "설 치 장 소", "라이선스 유형", "라이선스 보증기간"]
labels_right = ["고 객 번 호",    "",             "",              ""]
row_tops = [info_top + i * ROW_H for i in range(4)]

for i, (lbl, top) in enumerate(zip(labels_left, row_tops)):
    add_textbox(slide, COL1_L, top, COL1_W, ROW_H,
                text=lbl + " :", fontsize=10, bold=True,
                color=LABEL_COLOR, align=PP_ALIGN.LEFT)

# 우측 레이블 (1행만)
add_textbox(slide, COL3_L, row_tops[0], COL3_W, ROW_H,
            text="고 객 번 호 :", fontsize=10, bold=True,
            color=LABEL_COLOR, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════
#  Placeholder (app.py 가 읽는 idx 들)
#  idx 10: 고객명   idx 11: 설치장소   idx 12: 보증기간
#  idx 14: 고객번호   idx 18: 년   idx 19: 월   idx 20: 일
#  idx 21: 라이선스 유형
# ═══════════════════════════════════════════════════
PH_COLOR = RGBColor(0x22, 0x22, 0x22)

# idx 10: 고객명
add_placeholder_shape(slide, 10,
    int(COL2_L), int(row_tops[0]), int(COL2_W), int(ROW_H),
    default_text="고객명", fontsize=10, color=PH_COLOR)

# idx 11: 설치장소
add_placeholder_shape(slide, 11,
    int(COL2_L), int(row_tops[1]), int(COL2_W + COL4_W + Cm(0.5)), int(ROW_H),
    default_text="설치장소", fontsize=10, color=PH_COLOR)

# idx 21: 라이선스 유형
add_placeholder_shape(slide, 21,
    int(COL2_L), int(row_tops[2]), int(COL2_W + COL4_W + Cm(0.5)), int(ROW_H),
    default_text="라이선스 유형", fontsize=10, color=PH_COLOR)

# idx 12: 보증기간
add_placeholder_shape(slide, 12,
    int(COL2_L), int(row_tops[3]), int(COL2_W + COL4_W + Cm(0.5)), int(ROW_H),
    default_text="보증기간", fontsize=10, color=PH_COLOR)

# idx 14: 고객번호 (1행 우측)
add_placeholder_shape(slide, 14,
    int(COL4_L), int(row_tops[0]), int(COL4_W), int(ROW_H),
    default_text="고객번호", fontsize=10, color=PH_COLOR)


# ═══════════════════════════════════════════════════
#  구분선
# ═══════════════════════════════════════════════════
add_rect(slide, Cm(1.5), Cm(6.9), W - Cm(3.0), Cm(0.06), NAVY)


# ═══════════════════════════════════════════════════
#  라이선스 테이블
# ═══════════════════════════════════════════════════
tbl_top   = Cm(7.2)
tbl_left  = Cm(1.5)
tbl_w     = W - Cm(3.0)
tbl_h     = Cm(5.5)

rows_n = 4   # 헤더 1 + 데이터행 3 (app.py 가 동적으로 행을 추가/복사함)
cols_n = 3

tbl = slide.shapes.add_table(rows_n, cols_n, tbl_left, tbl_top, tbl_w, tbl_h).table

# 열 너비 비율: No(7%) / Software(78%) / QTY(15%)
total_emu = int(tbl_w)
tbl.columns[0].width = int(total_emu * 0.07)
tbl.columns[1].width = int(total_emu * 0.78)
tbl.columns[2].width = int(total_emu * 0.15)

# 행 높이
from pptx.util import Cm as C
for i in range(rows_n):
    tbl.rows[i].height = int(C(0.9) if i == 0 else C(0.85))

# 헤더
headers = ["No", "Software", "QTY"]
for c, h in enumerate(headers):
    set_cell(tbl.cell(0, c), h, fontsize=11, bold=True,
             fg=WHITE, bg=NAVY, align=PP_ALIGN.CENTER)

# 데이터행 (샘플)
sample = [
    ("1", "Ansys Mechanical Enterprise", "5 task(s)"),
    ("2", "Ansys Fluent", "2 task(s)"),
    ("3", "", ""),
]
for r, (no, sw, qty) in enumerate(sample, start=1):
    row_bg = LGRAY if r % 2 == 0 else WHITE
    set_cell(tbl.cell(r, 0), no,  fontsize=10, bg=row_bg, align=PP_ALIGN.CENTER)
    set_cell(tbl.cell(r, 1), sw,  fontsize=10, bg=row_bg, align=PP_ALIGN.LEFT)
    set_cell(tbl.cell(r, 2), qty, fontsize=10, bg=row_bg, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
#  발행일 (idx 18: 년 / 19: 월 / 20: 일)
# ═══════════════════════════════════════════════════
date_top = tbl_top + tbl_h + Cm(0.6)
date_total_w = Cm(9)
date_left = (W - date_total_w) // 2

# "____년   __월   __일"
add_textbox(slide, date_left, date_top, Cm(1.5), Cm(0.7),
            text="", fontsize=12)   # 여백

# 년 placeholder
add_placeholder_shape(slide, 18,
    int(date_left), int(date_top), int(Cm(2)), int(Cm(0.7)),
    default_text="2024", fontsize=12, bold=False,
    color=PH_COLOR, align=PP_ALIGN.RIGHT)

add_textbox(slide, date_left + Cm(2), date_top, Cm(0.7), Cm(0.7),
            text="년", fontsize=12, color=GRAY)

# 월 placeholder
add_placeholder_shape(slide, 19,
    int(date_left + Cm(2.8)), int(date_top), int(Cm(1.2)), int(Cm(0.7)),
    default_text="01", fontsize=12, bold=False,
    color=PH_COLOR, align=PP_ALIGN.RIGHT)

add_textbox(slide, date_left + Cm(4.1), date_top, Cm(0.5), Cm(0.7),
            text="월", fontsize=12, color=GRAY)

# 일 placeholder
add_placeholder_shape(slide, 20,
    int(date_left + Cm(4.8)), int(date_top), int(Cm(1.2)), int(Cm(0.7)),
    default_text="01", fontsize=12, bold=False,
    color=PH_COLOR, align=PP_ALIGN.RIGHT)

add_textbox(slide, date_left + Cm(6.1), date_top, Cm(0.5), Cm(0.7),
            text="일", fontsize=12, color=GRAY)


# ═══════════════════════════════════════════════════
#  하단 서명 영역
# ═══════════════════════════════════════════════════
sign_top = date_top + Cm(1.2)

add_textbox(slide, 0, sign_top, W, Cm(0.9),
            text="(주) 태성에스엔이  대표이사",
            fontsize=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_textbox(slide, 0, sign_top + Cm(0.9), W, Cm(0.6),
            text="TSNE CO., LTD.",
            fontsize=9, bold=False, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
#  저장
# ═══════════════════════════════════════════════════
out_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "라이선스_확인서_템플릿.pptx"
)
prs.save(out_path)
print(f"[OK] 템플릿 생성 완료: {out_path}")
print()
print("app.py 와 같은 폴더에 저장되었습니다.")
print("EXE 빌드 후에는 dist\\AnsysLicenseTool\\ 에도 복사해 주세요.")
