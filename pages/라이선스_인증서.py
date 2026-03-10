import streamlit as st
import re
import os
import io
import pandas as pd
from datetime import date
from copy import deepcopy
import sys
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))
from utils.theme import init_theme, render_theme_css, render_theme_toggle, muted_color

# ─── 테마 초기화 ────────────────────────────────────────────────
init_theme()

# ─── 커스텀 CSS ──────────────────────────────────────────────────
render_theme_css()

# ─── 테마 토글 버튼 ──────────────────────────────────────────────
render_theme_toggle()

# ─── 타이틀 ─────────────────────────────────────────────────────
st.markdown('<div class="page-tag">Ansys License Tool</div>', unsafe_allow_html=True)
st.markdown('<div class="page-title">라이선스 확인서 <em>자동 생성기</em></div>', unsafe_allow_html=True)
st.markdown('<div class="page-sub">라이선스 .txt 파일을 업로드하면 제품 목록을 자동 추출하고<br>공식 확인서 PDF 및 PPT를 즉시 다운로드할 수 있습니다.</div>', unsafe_allow_html=True)

# ─── STEP 1: 파일 업로드 ────────────────────────────────────────
st.markdown('<div class="section-label"><span class="sn">01</span> 라이선스 파일 업로드</div>', unsafe_allow_html=True)
uploaded_file = st.file_uploader(
    "라이선스 .txt 파일을 여기에 드래그하거나 클릭하여 업로드하세요",
    type=["txt"], label_visibility="collapsed"
)

df_licenses = pd.DataFrame()

# ─── STEP 2: 파싱 & 테이블 ─────────────────────────────────────
st.markdown('<div class="section-label"><span class="sn">02</span> 추출된 라이선스 목록</div>', unsafe_allow_html=True)

extracted_customer_no = None
extracted_warranty_end = None

if uploaded_file:
    content = uploaded_file.read().decode("utf-8", errors="ignore")
    pattern = re.compile(
        r'#?\s*(\d+)\.\s+([\w\s\(\)\/\-]+?):\s+(\d+)\s+task\(s\).*?expiring\s+([\d\-a-zA-Z]+).*?Customer\s*#\s*(\d+)',
        re.IGNORECASE | re.DOTALL
    )
    matches = pattern.findall(content)
    if matches:
        df_licenses = pd.DataFrame(matches, columns=["No", "Software (제품명)", "QTY (수량)", "ExpireDate", "CustomerNo"])
        if not df_licenses.empty:
            extracted_customer_no = df_licenses.iloc[-1]["CustomerNo"]
            extracted_warranty_end = df_licenses.iloc[-1]["ExpireDate"]
        df_display = df_licenses[["No", "Software (제품명)", "QTY (수량)"]]
        st.dataframe(df_display, use_container_width=True, hide_index=True)
        st.success(f"✅ 총 {len(df_licenses)}개 라이선스 항목이 감지되었습니다.")
        df_licenses = df_display
    else:
        st.warning("⚠️ 파일에서 라이선스 패턴을 찾지 못했습니다.")
        with st.expander("파일 내용 미리보기"):
            st.text(content[:2000])
else:
    st.info("① 위에서 파일을 업로드하면 여기에 라이선스 목록이 표시됩니다.")

# ─── STEP 3: 확인서 폼 ──────────────────────────────────────────
st.markdown('<div class="section-label"><span class="sn">03</span> 확인서 정보 입력</div>', unsafe_allow_html=True)

with st.form("cert_form"):
    col1, col2 = st.columns(2)
    with col1:
        customer_name    = st.text_input("고 객 명",    placeholder="예) 한국 컴퍼니")
        customer_number  = st.text_input("고 객 번 호", value=extracted_customer_no or "", placeholder="예) 1213401")
    with col2:
        install_location = st.text_input("설 치 장 소",    placeholder="예) 서울시 강남구 테헤란로 123")
        warranty_start   = st.date_input("라이선스 보증기간 시작", value=date.today())

    warranty_end_default = date.today()
    if extracted_warranty_end:
        try:
            warranty_end_default = pd.to_datetime(extracted_warranty_end).date()
        except Exception:
            pass

    warranty_end = st.date_input("라이선스 보증기간 끝", value=warranty_end_default)

    license_type = st.selectbox(
        "라이선스 유형",
        ["Permanent License / LAN", "Lease License / WAN",
         "Maintenance License / LAN", "Permanent License / Regional WAN"]
    )
    issue_date = st.date_input("발행 일자", value=date.today())
    submitted  = st.form_submit_button("📋  확인서 생성 (PPT)", use_container_width=True)




# ══════════════════════════════════════════════════════════
#  PPTX 생성 — 원본 템플릿 기반
# ══════════════════════════════════════════════════════════
def create_pptx_from_template(customer, customer_no, location, lic_type,
                               warranty, iss_date, df, template_bytes):
    from pptx import Presentation

    prs = Presentation(io.BytesIO(template_bytes))
    slide = prs.slides[0]

    ph_map = {
        10: customer,
        11: location,
        12: warranty,
        14: customer_no,
        18: str(iss_date.year),
        19: f"{iss_date.month:02d}",
        20: f"{iss_date.day:02d}",
        21: lic_type,
    }

    for shape in slide.shapes:
        try:
            ph = shape.placeholder_format
            if ph and ph.idx in ph_map:
                tf = shape.text_frame
                para = tf.paragraphs[0]
                if para.runs:
                    para.runs[0].text = ph_map[ph.idx]
                else:
                    para.text = ph_map[ph.idx]
        except Exception:
            pass

    if not df.empty:
        items = [(str(r["No"]), str(r["Software (제품명)"]),
                  str(r["QTY (수량)"]) + " task(s)")
                 for _, r in df.iterrows()]
    else:
        items = [("1", "라이선스 정보 없음", "-")]

    n_items = len(items)
    cell_pt = 10 if n_items <= 7 else 9 if n_items <= 12 else 8

    from pptx.util import Pt as _Pt, Cm as _Cm
    max_data_h = _Cm(5.1)

    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            while len(tbl.rows) < len(items):
                last_tr = tbl.rows[len(tbl.rows) - 1]._tr
                new_tr = deepcopy(last_tr)
                tbl._tbl.append(new_tr)

            if len(tbl.rows) * _Cm(0.85) > max_data_h:
                new_h = int(max_data_h / len(tbl.rows))
                for i in range(len(tbl.rows)):
                    tbl.rows[i].height = new_h

            for r_idx, (no, sw, qty) in enumerate(items):
                row = tbl.rows[r_idx]
                for c_idx, val in enumerate([no, sw, qty]):
                    cell = row.cells[c_idx]
                    tf = cell.text_frame
                    para = tf.paragraphs[0]
                    if para.runs:
                        para.runs[0].text = val
                        para.runs[0].font.size = _Pt(cell_pt)
                    else:
                        para.text = val
            break

    sld_id_list = prs.slides._sldIdLst
    sld_ids = list(sld_id_list)
    for i in range(1, len(sld_ids)):
        sld_id_list.remove(sld_ids[i])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─── 생성 버튼 처리 ─────────────────────────────────────────────
if submitted:
    if not customer_name or not install_location:
        st.error("❌ 고객명과 설치 장소를 모두 입력해 주세요.")
    else:
        template_bytes = None
        # 템플릿 파일을 pages/ 기준이 아닌 프로젝트 루트 기준으로 찾기
        root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        template_candidates = [
            os.path.join(root_dir, "라이선스_확인서_템플릿.pptx"),
            "라이선스_확인서_템플릿.pptx",
        ]
        for tp in template_candidates:
            if os.path.exists(tp):
                with open(tp, "rb") as f:
                    template_bytes = f.read()
                break

        with st.spinner("확인서 생성 중..."):
            fname_base = f"Ansys_라이선스확인서_{customer_name}_{issue_date.strftime('%Y%m%d')}"
            warranty_period_formatted = (
                f"{warranty_start.year}. {warranty_start.month:02d}. {warranty_start.day:02d}"
                f" ~ {warranty_end.year}. {warranty_end.month:02d}. {warranty_end.day:02d}"
            )

            _loc = install_location or "-"
            pptx_buf  = None
            pptx_err  = None
            if template_bytes:
                try:
                    pptx_buf = create_pptx_from_template(
                        customer_name, customer_number or "-",
                        _loc, license_type,
                        warranty_period_formatted,
                        issue_date, df_licenses,
                        template_bytes
                    )
                except Exception as e:
                    pptx_err = str(e)
            else:
                pptx_err = "템플릿 파일(라이선스_확인서_템플릿.pptx)이 앱 폴더에 없습니다."

        st.success("✅ 확인서가 생성되었습니다! 아래 버튼으로 다운로드하세요.")

        if pptx_buf:
            st.download_button(
                label="⬇️  PPT 템플릿 다운로드",
                data=pptx_buf,
                file_name=fname_base + ".pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
        else:
            st.error(f"PPT 오류: {pptx_err}")

st.markdown("---")
st.markdown(
    f"<center style='color:{muted_color()}; font-size:.75rem; font-family:monospace;'>"
    "Ansys License Certificate Generator · Navy × Yellow Edition"
    "</center>",
    unsafe_allow_html=True
)
